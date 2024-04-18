from langchain.chat_models import ChatOpenAI
from langchain.tools import tool
from langchain.prompts import ChatPromptTemplate
from langchain.schema.output_parser import StrOutputParser
from langchain.output_parsers.json import SimpleJsonOutputParser  # JsonOutputParser 임포트
from langchain.schema.runnable import RunnablePassthrough

class PresentationTools():
    @tool("Generate Presentation Slides")
    def generate_presentation_slides(topic):
        """
        Generates presentation slides for a given topic using OpenAI's GPT-3.5-turbo model.
        """
        prompt = f"Generate a 10 slide presentation for the topic. Produce 50 to 60 words per slide. {topic}. Each slide should have a header and content. The final slide should be a list of discussion questions. Return as JSON."
        chat_model = ChatOpenAI(model="gpt-3.5-turbo")
        response = chat_model.run(prompt)
        return response
    
from langchain.chains import Chain
from pptx import Presentation

def create_presentation_chain():
    chain = Chain(steps=[
        {
            "name": "Generate Presentation Slides",
            "tool": "PresentationTools.generate_presentation_slides",
            "inputs": {"topic": "{{ topic }}"},
            "outputs": {"slides": "{{ result }}"},
        },
        {
            "name": "Create PowerPoint Presentation",
            "runnable": lambda slides: create_powerpoint(slides),
            "inputs": {"slides": "{{ slides }}"},
        },
    ])
    return chain

def create_powerpoint(slide_data):
    prs = Presentation()
    for slide in slide_data:
        slide_layout = prs.slide_layouts[1]
        new_slide = prs.slides.add_slide(slide_layout)
        if slide['header']:
            title = new_slide.shapes.title
            title.text = slide['header']
        if slide['content']:
            shapes = new_slide.shapes
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.text = slide['content']
            tf.fit_text(font_family="Calibri", max_size=18, bold=True)
    prs.save("output.pptx")

chain = create_presentation_chain()
chain.invoke({"topic": "Your Presentation Topic"})

