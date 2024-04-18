import openai, json
from pptx import Presentation
from typing import Union
#openai.api_key = "YOUR_API_KEY_HERE"

#presentation_title = input("What do you want to make a presentation about?")

from langchain.prompts import ChatPromptTemplate
from langchain.schema.output_parser import StrOutputParser
from langchain.output_parsers.json import SimpleJsonOutputParser  # JsonOutputParser 임포트
from langchain.schema.runnable import RunnablePassthrough
from langchain.chat_models import ChatOpenAI


prompt_template = ChatPromptTemplate.from_template(
    """
    Generate a 10 slide presentation for the topic. 
    Produce 50 to 60 words per slide. 
    topic: {topic}. 
    Each slide should have a header and content. The final slide should be a list of discussion questions. 
    Return as JSON.

    {{"slides":[
        {{
            "header": "header of the slide",
            "content": "content of the slide"
        }}

    ]}}
    """
    # template="",
    # variables={"topic": presentation_title},
    # response_parser="json",
    # response_parser_args={"path": "slides"}
)



from langchain.tools import tool

from langchain.pydantic_v1 import BaseModel, Field
from typing import List
import json

class SlideModel(BaseModel):
    header: str = Field(..., description="The title or main heading of the slide.")
    content: str = Field(..., description="The main body text of the slide.")

class PresentationModel(BaseModel):
    slides: List[SlideModel] = Field(..., description="A list of slides that make up the presentation.")

    
# Convert response to PresentationModel
#presentation_model = PresentationModel(**response)

class PowerpointTools():
    @tool("slide-generator", return_direct=True)
    def generate_slide(slides: str):
        """
        This method takes string representing JSON presentation: ```{'slides':[{'header':'header of the slide', content:'content of the slide'}]}```.
        Each slide in the presentation is created based on the header and content provided in the presentation data.
        The final PowerPoint file is returned, ready for saving or further manipulation.
        """
        prs = Presentation()

        # Check if slides input is a string, and if so, parse it as JSON into SlideModel
        if isinstance(slides, str):
            slides_dict = json.loads(slides)
            slides = PresentationModel(**slides_dict).slides
        # If slides is already a SlideModel instance, use it directly
        elif isinstance(slides, SlideModel):
            slides = [slides]  # Convert single SlideModel instance to list for iteration

        for slide in slides:
            slide_layout = prs.slide_layouts[1]
            new_slide = prs.slides.add_slide(slide_layout)
            
            if slide.header:
                title = new_slide.shapes.title
                title.text = slide.header
            
            if slide.content:
                shapes = new_slide.shapes
                body_shape = shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.text = slide.content
                tf.fit_text(font_family="Arial", max_size=18, bold=True)
                
        import uuid
        filename = "output/" + str(uuid.uuid4()) + ".pptx"
        prs.save(filename)

        return filename

# prs = PowerpointTool.generate_slide(response)
# prs.save("output.pptx")


# model = ChatOpenAI(model="gpt-3.5-turbo")

# chain = (RunnablePassthrough() | prompt_template | model | StrOutputParser())

#response = chain.invoke({"topic": "클라우드 네이티브 앱을 정부에 적용할 제안서"})


if __name__ == "__main__":
#print(PowerpointTools.generate_slide("{slides:[{header:'title', content: 'content'}]}"))
    tools = [PowerpointTools.generate_slide]

    from langchain import hub

    # Get the prompt to use - you can modify this!
    prompt = hub.pull("hwchase17/openai-functions-agent")
    prompt.messages

    from langchain.agents import create_openai_functions_agent
    
    llm = ChatOpenAI(model="gpt-3.5-turbo")

    agent = create_openai_functions_agent(llm, tools, prompt)

    from langchain.agents import AgentExecutor

    agent_executor = AgentExecutor(agent=agent, tools=tools, verbose=True)


    from langchain_core.messages import AIMessage, HumanMessage


    agent_executor.invoke(
        {
            "chat_history": [
                HumanMessage(content="please generate a slides for presenting about cloud-native applications within 5 slides"),
                AIMessage(content="Sure! I can help you generate slides for presenting about cloud-native applications. Could you please provide me with the content for each slide?"),
            ],
            "input": "1. what is cloud-native applications, 2. benefits of CNAs, 3. architecture of CNAs, 4. transformation strategies, 5. conclusion",
        }
    )



    

